import streamlit as st
import pandas as pd
from datetime import date, timedelta
from pathlib import Path
import time
import json
import re
import os
import io
import subprocess
import shutil
import secrets
from typing import List, Dict, Any, Optional

# -----------------------------
# Keep-warm ping (must be first)
# -----------------------------
qp = st.query_params
if qp.get("ping") == "1":
    st.write("ok")
    st.stop()

# -----------------------------
# Optional libraries
# -----------------------------
try:
    from pypdf import PdfReader
    PDF_TEXT_OK = True
except Exception:
    PDF_TEXT_OK = False

try:
    import fitz  # PyMuPDF
    PDF_RENDER_OK = True
except Exception:
    PDF_RENDER_OK = False

try:
    from pptx import Presentation
    PPTX_OK = True
except Exception:
    PPTX_OK = False

try:
    import pytesseract
    from PIL import Image
    OCR_OK = True
except Exception:
    OCR_OK = False

try:
    from groq import Groq
    GROQ_OK = True
except Exception:
    GROQ_OK = False


# -----------------------------
# Page config
# -----------------------------
st.set_page_config(page_title="Smart Study Planner", page_icon="📚", layout="wide")

# -----------------------------
# Helpers
# -----------------------------
def sanitize_name(name: str) -> str:
    return name.replace("/", "_").replace("\\", "_")

def now_str() -> str:
    return pd.Timestamp.now().strftime("%Y-%m-%d %H:%M")

def normalize_spaces(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "")).strip()

def parse_json_loose(raw: str):
    raw = (raw or "").strip()
    try:
        return json.loads(raw)
    except Exception:
        pass
    start = raw.find("{")
    end = raw.rfind("}")
    if start != -1 and end != -1 and end > start:
        return json.loads(raw[start:end + 1])
    raise ValueError("No JSON object found")

def nice_date(d) -> str:
    try:
        return pd.to_datetime(d).strftime("%b %d, %Y")
    except Exception:
        return ""

def default_df() -> pd.DataFrame:
    t = date.today()
    return pd.DataFrame(
        {
            "Subject": ["Physics", "Gen Math", "Biology"],
            "Difficulty (1-5)": [5, 4, 3],
            "Exam Date": [t + timedelta(days=10), t + timedelta(days=7), t + timedelta(days=14)],
            "Minutes Done (this week)": [0, 0, 0],
        }
    )


# -----------------------------
# Workspace (private per key) - SIDEBAR
# -----------------------------
with st.sidebar:
    st.markdown("## 🔐 Workspace")
    if "workspace_key" not in st.session_state:
        st.session_state.workspace_key = ""

    if not st.session_state.workspace_key:
        key_in = st.text_input("Workspace key", type="password", placeholder="Enter your key")
        c1, c2 = st.columns(2)
        with c1:
            if st.button("Create"):
                st.session_state.workspace_key = secrets.token_urlsafe(16)
                st.success("Created — copy & save this key!")
                st.code(st.session_state.workspace_key)
        with c2:
            if st.button("Use"):
                if key_in.strip():
                    st.session_state.workspace_key = key_in.strip()
                    st.rerun()
                else:
                    st.warning("Enter a key first.")

        st.info("This key protects your files. If you lose it, you lose access to your saved workspace.")
        st.stop()
    else:
        st.success("Workspace active")
        if st.button("Switch workspace"):
            st.session_state.workspace_key = ""
            st.rerun()

    st.divider()
    st.markdown("## ⚙️ Settings")
    minutes_per_day = st.number_input("Minutes/day", 30, 600, 180, 10)
    days_to_plan = st.slider("Plan length (days)", 3, 14, 7)

# storage key
user_code = st.session_state.workspace_key
safe_code = re.sub(r"[^a-zA-Z0-9_-]", "", user_code)[:40] or "anon"

DATA_PATH = Path(f"data_{safe_code}.csv")
FILES_DIR = Path(f"files_{safe_code}")
FILES_DIR.mkdir(exist_ok=True)
META_PATH = FILES_DIR / "files_meta.json"


# -----------------------------
# Data IO
# -----------------------------
def load_df() -> pd.DataFrame:
    if DATA_PATH.exists():
        try:
            df = pd.read_csv(DATA_PATH)
            df["Exam Date"] = pd.to_datetime(df["Exam Date"], errors="coerce").dt.date
            return df
        except Exception:
            return default_df()
    return default_df()

def save_df(df: pd.DataFrame) -> None:
    out = df.copy()
    out["Exam Date"] = pd.to_datetime(out["Exam Date"], errors="coerce").dt.strftime("%Y-%m-%d")
    out.to_csv(DATA_PATH, index=False)

def load_meta() -> Dict[str, Any]:
    if META_PATH.exists():
        try:
            return json.loads(META_PATH.read_text(encoding="utf-8"))
        except Exception:
            return {"files": []}
    return {"files": []}

def save_meta(meta: Dict[str, Any]) -> None:
    META_PATH.write_text(json.dumps(meta, indent=2), encoding="utf-8")


# -----------------------------
# Extraction / conversion
# -----------------------------
def extract_text_from_pdf_text(p: Path, max_pages: int = 12) -> str:
    if not (PDF_TEXT_OK and p.suffix.lower() == ".pdf"):
        return ""
    try:
        reader = PdfReader(str(p))
        parts = []
        for page in reader.pages[:max_pages]:
            parts.append(page.extract_text() or "")
        return "\n".join(parts).strip()
    except Exception:
        return ""

def ocr_pdf_with_pymupdf(p: Path, max_pages: int = 6, zoom: float = 2.0) -> str:
    if not (PDF_RENDER_OK and OCR_OK and p.suffix.lower() == ".pdf"):
        return ""
    try:
        doc = fitz.open(str(p))
        pages = min(max_pages, doc.page_count)
        parts = []
        for i in range(pages):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            text = (pytesseract.image_to_string(img) or "").strip()
            if text:
                parts.append(f"[Page {i+1}]\n{text}")
        doc.close()
        return "\n\n".join(parts).strip()
    except Exception:
        return ""

def extract_text_from_pptx(p: Path, max_slides: int = 60) -> str:
    if not (PPTX_OK and p.suffix.lower() == ".pptx"):
        return ""
    try:
        prs = Presentation(str(p))
        out = []
        for i, slide in enumerate(prs.slides):
            if i >= max_slides:
                break
            chunks = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = (shape.text or "").strip()
                    if t:
                        chunks.append(t)
            if chunks:
                out.append(f"[Slide {i+1}]\n" + "\n".join(chunks))
        return "\n\n".join(out).strip()
    except Exception:
        return ""

def extract_text_smart(p: Path) -> str:
    suf = p.suffix.lower()
    if suf == ".txt":
        try:
            return (p.read_text(errors="ignore") or "").strip()
        except Exception:
            return ""
    if suf == ".pdf":
        txt = extract_text_from_pdf_text(p, max_pages=12)
        if len(normalize_spaces(txt)) < 200:
            ocr_txt = ocr_pdf_with_pymupdf(p, max_pages=6, zoom=2.0)
            if ocr_txt:
                return ocr_txt
        return txt
    if suf == ".pptx":
        return extract_text_from_pptx(p, max_slides=60)
    return ""

def pptx_to_pdf(pptx_path: Path) -> Optional[Path]:
    out_dir = pptx_path.parent
    out_pdf = out_dir / f"{pptx_path.stem}.pdf"

    try:
        if out_pdf.exists() and out_pdf.stat().st_mtime >= pptx_path.stat().st_mtime:
            return out_pdf
    except Exception:
        pass

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None

    try:
        subprocess.run(
            [
                soffice,
                "--headless",
                "--nologo",
                "--nofirststartwizard",
                "--convert-to",
                "pdf",
                "--outdir",
                str(out_dir),
                str(pptx_path),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        return out_pdf if out_pdf.exists() else None
    except Exception:
        return None


# -----------------------------
# Session init
# -----------------------------
if "active_user_code" not in st.session_state or st.session_state.active_user_code != safe_code:
    st.session_state.active_user_code = safe_code
    st.session_state.subjects = load_df()

if "plan_ready" not in st.session_state:
    st.session_state.plan_ready = False

if "timer_running" not in st.session_state:
    st.session_state.timer_running = False
if "timer_subject" not in st.session_state:
    st.session_state.timer_subject = ""
if "timer_minutes" not in st.session_state:
    st.session_state.timer_minutes = 25
if "timer_start" not in st.session_state:
    st.session_state.timer_start = 0.0
if "timer_end" not in st.session_state:
    st.session_state.timer_end = 0.0

if "preview_path" not in st.session_state:
    st.session_state.preview_path = ""
if "preview_name" not in st.session_state:
    st.session_state.preview_name = ""
if "preview_subject" not in st.session_state:
    st.session_state.preview_subject = ""


# -----------------------------
# Header
# -----------------------------
st.markdown("# 📚 Smart Study Planner")
st.caption("From the works of STEM 12 A")

tabs = st.tabs(["📌 Subjects", "📁 Files + Preview", "⏱ Timer", "🧠 Assessment", "📊 Dashboard"])

# -----------------------------
# TAB: Subjects
# -----------------------------
with tabs[0]:
    st.subheader("Subjects")
    edited_df = st.data_editor(
        st.session_state.subjects,
        num_rows="dynamic",
        use_container_width=True,
        column_config={
            "Exam Date": st.column_config.DateColumn("Exam Date", format="MMM DD, YYYY"),
            "Difficulty (1-5)": st.column_config.NumberColumn("Difficulty (1-5)", min_value=1, max_value=5, step=1),
            "Minutes Done (this week)": st.column_config.NumberColumn("Minutes Done (this week)", min_value=0, step=10),
        },
    )

    cA, cB = st.columns([1, 1])
    with cA:
        if st.button("💾 Save changes", key="save_subjects"):
            st.session_state.subjects = edited_df.copy()
            save_df(st.session_state.subjects)
            st.success("Saved!")
    with cB:
        if st.button("🗑 Reset to default", key="reset_subjects"):
            st.session_state.subjects = default_df()
            save_df(st.session_state.subjects)
            st.success("Reset done!")
            st.rerun()

    subjects_clean = (
        pd.Series(st.session_state.subjects["Subject"])
        .dropna()
        .astype(str)
        .str.strip()
    )
    subjects_clean = subjects_clean[subjects_clean != ""].tolist()

    st.divider()
    if st.button("✅ Update Plan", key="update_plan_btn"):
        st.session_state.plan_ready = True
        st.success("Plan updated. Check Dashboard tab.")

# refresh subjects_clean for other tabs
subjects_clean = (
    pd.Series(st.session_state.subjects["Subject"])
    .dropna()
    .astype(str)
    .str.strip()
)
subjects_clean = subjects_clean[subjects_clean != ""].tolist()

# -----------------------------
# TAB: Files + Preview
# -----------------------------
with tabs[1]:
    st.subheader("Study Files")
    meta = load_meta()

    with st.expander("Upload study files", expanded=True):
        if not subjects_clean:
            st.warning("Add at least 1 subject first so files can be tagged.")
        else:
            tag_subject = st.selectbox("Tag subject", subjects_clean, key="tag_subject_upload")
            uploads = st.file_uploader(
                "Upload (TXT / PDF / PPTX / images)",
                type=["txt", "pdf", "pptx", "png", "jpg", "jpeg"],
                accept_multiple_files=True,
                key="uploader",
            )
            if st.button("⬆️ Save uploaded files", key="save_uploads"):
                if not uploads:
                    st.warning("No files selected.")
                else:
                    subj_dir = FILES_DIR / tag_subject
                    subj_dir.mkdir(exist_ok=True)
                    added = 0
                    for f in uploads:
                        fname = sanitize_name(f.name)
                        path = subj_dir / fname
                        path.write_bytes(f.getbuffer())
                        meta["files"].append(
                            {"subject": tag_subject, "name": fname, "path": str(path), "uploaded_at": now_str()}
                        )
                        added += 1
                    save_meta(meta)
                    st.success(f"Saved {added} file(s) under {tag_subject}.")
                    st.rerun()

    st.divider()

    if not meta.get("files"):
        st.info("No study files yet.")
    else:
        files_by_subject: Dict[str, List[Dict[str, Any]]] = {}
        for item in meta["files"]:
            files_by_subject.setdefault(item["subject"], []).append(item)

        for subj in sorted(files_by_subject.keys(), key=lambda x: x.lower()):
            with st.expander(f"📚 {subj}", expanded=False):
                items = sorted(files_by_subject[subj], key=lambda x: x["name"].lower())
                for item in items:
                    p = Path(item["path"])
                    if not p.exists():
                        st.warning(f"Missing file: {item['name']}")
                        continue

                    col1, col2, col3, col4 = st.columns([7, 1.5, 1.5, 1.5])
                    with col1:
                        st.write(f"**{item['name']}**  \nUploaded: {item.get('uploaded_at','')}")
                    with col2:
                        st.download_button(
                            "Download",
                            data=p.read_bytes(),
                            file_name=item["name"],
                            mime="application/octet-stream",
                            key=f"dl_{subj}_{item['name']}",
                        )
                    with col3:
                        if st.button("Delete", key=f"del_{subj}_{item['name']}"):
                            p.unlink(missing_ok=True)
                            meta["files"] = [
                                x for x in meta["files"]
                                if not (x["subject"] == subj and x["name"] == item["name"])
                            ]
                            save_meta(meta)
                            if st.session_state.preview_path == str(p):
                                st.session_state.preview_path = ""
                                st.session_state.preview_name = ""
                                st.session_state.preview_subject = ""
                            st.rerun()
                    with col4:
                        if st.button("Preview", key=f"prev_{subj}_{item['name']}"):
                            st.session_state.preview_path = str(p)
                            st.session_state.preview_name = item["name"]
                            st.session_state.preview_subject = subj
                            st.rerun()

    # Preview area (kept contained)
    if st.session_state.preview_path:
        p = Path(st.session_state.preview_path)
        if p.exists():
            st.divider()
            st.subheader(f"🔎 Preview: {st.session_state.preview_subject} — {st.session_state.preview_name}")

            if st.button("❌ Close preview", key="close_preview"):
                st.session_state.preview_path = ""
                st.session_state.preview_name = ""
                st.session_state.preview_subject = ""
                st.rerun()

            tabv, tabt = st.tabs(["📄 View file", "📝 Extracted text"])

            with tabv:
                suf = p.suffix.lower()

                if suf == ".pdf":
                    st.download_button("⬇️ Download PDF", data=p.read_bytes(), file_name=p.name, mime="application/pdf")

                    if not PDF_RENDER_OK:
                        st.warning("PDF viewer needs pymupdf.")
                    else:
                        doc = fitz.open(str(p))
                        total = doc.page_count

                        mode = st.selectbox("View mode", ["Single page", "Continuous (paged)"], key=f"pdf_mode_{p.name}")
                        zoom = st.slider("Zoom", 1.0, 3.0, 1.8, 0.1, key=f"pdf_zoom_{p.name}")
                        fit = st.toggle("Fit to screen", value=True, key=f"pdf_fit_{p.name}")

                        if mode == "Single page":
                            page_num = st.number_input("Page", 1, total, 1, 1, key=f"pdf_page_{p.name}")
                            page = doc.load_page(int(page_num) - 1)
                            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                            b = pix.tobytes("png")
                            if fit:
                                st.image(b, use_container_width=True, caption=f"Page {page_num}/{total}")
                            else:
                                st.image(b, width=int(900 * zoom), caption=f"Page {page_num}/{total}")
                        else:
                            per_batch = st.select_slider("Pages per batch", [2,3,4,5,6,8,10,15,20], value=6, key=f"pdf_batch_{p.name}")
                            start = st.number_input("Start page", 1, total, 1, 1, key=f"pdf_start_{p.name}")
                            end = min(total, int(start) + int(per_batch) - 1)

                            for i in range(int(start) - 1, end):
                                page = doc.load_page(i)
                                pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                                b = pix.tobytes("png")
                                with st.expander(f"Page {i+1}", expanded=(i == int(start) - 1)):
                                    if fit:
                                        st.image(b, use_container_width=True)
                                    else:
                                        st.image(b, width=int(900 * zoom))
                        doc.close()

                elif suf == ".pptx":
                    st.download_button("⬇️ Download PPTX", data=p.read_bytes(), file_name=p.name,
                                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

                    if not PDF_RENDER_OK:
                        st.warning("Slide viewer needs pymupdf.")
                    else:
                        pdf_path = pptx_to_pdf(p)
                        if not pdf_path:
                            st.warning("PPTX→PDF conversion unavailable. Ensure LibreOffice is installed.")
                        else:
                            doc = fitz.open(str(pdf_path))
                            total = doc.page_count

                            mode = st.selectbox("View mode", ["Single slide", "Continuous (paged)"], key=f"pptx_mode_{p.name}")
                            zoom = st.slider("Zoom", 1.0, 3.0, 1.8, 0.1, key=f"pptx_zoom_{p.name}")
                            fit = st.toggle("Fit to screen", value=True, key=f"pptx_fit_{p.name}")

                            if mode == "Single slide":
                                slide_num = st.number_input("Slide", 1, total, 1, 1, key=f"pptx_slide_{p.name}")
                                page = doc.load_page(int(slide_num) - 1)
                                pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                                b = pix.tobytes("png")
                                if fit:
                                    st.image(b, use_container_width=True, caption=f"Slide {slide_num}/{total}")
                                else:
                                    st.image(b, width=int(900 * zoom), caption=f"Slide {slide_num}/{total}")
                            else:
                                per_batch = st.select_slider("Slides per batch", [2,3,4,5,6,8,10,15,20], value=6, key=f"pptx_batch_{p.name}")
                                start = st.number_input("Start slide", 1, total, 1, 1, key=f"pptx_start_{p.name}")
                                end = min(total, int(start) + int(per_batch) - 1)

                                for i in range(int(start) - 1, end):
                                    page = doc.load_page(i)
                                    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                                    b = pix.tobytes("png")
                                    with st.expander(f"Slide {i+1}", expanded=(i == int(start) - 1)):
                                        if fit:
                                            st.image(b, use_container_width=True)
                                        else:
                                            st.image(b, width=int(900 * zoom))
                            doc.close()

                elif suf == ".txt":
                    st.text_area("Text file", p.read_text(errors="ignore"), height=420)
                else:
                    st.image(str(p), use_container_width=True)

            with tabt:
                suf = p.suffix.lower()
                if suf in [".txt", ".pdf", ".pptx"]:
                    txt = extract_text_smart(p)
                    if not txt.strip():
                        st.warning("No text extracted.")
                    else:
                        st.text_area("Extracted text", txt, height=420)
                else:
                    st.info("Text extraction is available for TXT/PDF/PPTX only.")


# -----------------------------
# Compute plan (only if updated)
# -----------------------------
df = st.session_state.subjects.copy()
df = df.dropna(subset=["Subject"])
df["Subject"] = df["Subject"].astype(str).str.strip()
df = df[df["Subject"] != ""]
today = date.today()
weekly_minutes = int(minutes_per_day * 7)

if not df.empty:
    df["Difficulty (1-5)"] = pd.to_numeric(df["Difficulty (1-5)"], errors="coerce").fillna(0)
    df["Minutes Done (this week)"] = pd.to_numeric(df["Minutes Done (this week)"], errors="coerce").fillna(0)

    exam_ts = pd.to_datetime(df["Exam Date"], errors="coerce").fillna(pd.Timestamp(today))
    df["Exam Date"] = exam_ts.dt.date
    df["Days Left"] = exam_ts.dt.date.apply(lambda d: max(1, (d - today).days)).astype(int)
    df["Urgency"] = 10 / df["Days Left"]
    df["Priority"] = (df["Difficulty (1-5)"] * df["Urgency"]).fillna(0)

    total_priority = float(df["Priority"].sum())
    if total_priority <= 0:
        df["Minutes/Week (Suggested)"] = 0
    else:
        df["Minutes/Week (Suggested)"] = ((df["Priority"] / total_priority) * weekly_minutes).fillna(0).round().astype(int)

    df["Progress %"] = (df["Minutes Done (this week)"] / df["Minutes/Week (Suggested)"]).fillna(0).clip(0, 1)


# -----------------------------
# TAB: Timer
# -----------------------------
with tabs[2]:
    st.subheader("⏱ Focus Timer (auto-logs minutes)")

    if df.empty:
        st.info("Add subjects first.")
    else:
        subjects_list = df["Subject"].astype(str).tolist()

        if not st.session_state.timer_running:
            chosen = st.selectbox("Choose subject", subjects_list, key="timer_subject_select")
            mins = st.number_input("Minutes", 1, 120, int(st.session_state.timer_minutes))
            c1, c2 = st.columns(2)
            with c1:
                if st.button("▶ Start Timer"):
                    st.session_state.timer_running = True
                    st.session_state.timer_subject = str(chosen)
                    st.session_state.timer_minutes = int(mins)
                    st.session_state.timer_start = time.time()
                    st.session_state.timer_end = st.session_state.timer_start + int(mins) * 60
                    st.rerun()
            with c2:
                if st.button("🛑 Cancel"):
                    st.session_state.timer_running = False
                    st.rerun()
        else:
            remaining = int(st.session_state.timer_end - time.time())
            planned = int(st.session_state.timer_minutes)
            subj = st.session_state.timer_subject

            elapsed_minutes = int((time.time() - st.session_state.timer_start) // 60)
            elapsed_minutes = max(0, min(planned, elapsed_minutes))

            b1, b2 = st.columns(2)
            with b1:
                if st.button("⏹ Stop & Log elapsed"):
                    if elapsed_minutes > 0:
                        base = load_df()
                        base["Minutes Done (this week)"] = pd.to_numeric(base["Minutes Done (this week)"], errors="coerce").fillna(0)
                        mask = base["Subject"].astype(str) == str(subj)
                        if mask.any():
                            base.loc[mask, "Minutes Done (this week)"] += elapsed_minutes
                            save_df(base)
                            st.session_state.subjects = base
                    st.session_state.timer_running = False
                    st.rerun()

            with b2:
                if st.button("🛑 Cancel (no log)"):
                    st.session_state.timer_running = False
                    st.rerun()

            if remaining <= 0:
                base = load_df()
                base["Minutes Done (this week)"] = pd.to_numeric(base["Minutes Done (this week)"], errors="coerce").fillna(0)
                mask = base["Subject"].astype(str) == str(subj)
                if mask.any():
                    base.loc[mask, "Minutes Done (this week)"] += planned
                    save_df(base)
                    st.session_state.subjects = base
                st.session_state.timer_running = False
                st.success(f"✅ Logged {planned} minutes to {subj}")
                st.rerun()
            else:
                st.info(f"Studying **{subj}** — {remaining//60:02d}:{remaining%60:02d} (elapsed: {elapsed_minutes} min)")
                time.sleep(1)
                st.rerun()


# -----------------------------
# TAB: Assessment
# -----------------------------
with tabs[3]:
    st.subheader("🧠 Assessment (Answer first → AI checks later)")

    groq_key = st.secrets.get("GROQ_API_KEY") or os.environ.get("GROQ_API_KEY")
    ai_available = bool(groq_key) and GROQ_OK

    if not ai_available:
        st.warning("Assessment AI is OFF. Add GROQ_API_KEY in Secrets and include 'groq' in requirements.txt.")
    elif not subjects_clean:
        st.info("Add subjects and upload notes first.")
    else:
        client = Groq(api_key=groq_key)

        assess_subject = st.selectbox("Subject", subjects_clean, key="assess_subject_groq")

        meta = load_meta()
        subject_paths: List[Path] = []
        for item in meta.get("files", []):
            if item.get("subject") == assess_subject:
                pf = Path(item.get("path", ""))
                if pf.exists() and pf.suffix.lower() in [".txt", ".pdf", ".pptx"]:
                    subject_paths.append(pf)

        if not subject_paths:
            st.warning("Upload at least one TXT/PDF/PPTX notes file for this subject.")
        else:
            chosen_file = st.selectbox("Notes file", [p.name for p in subject_paths], key="assess_file_groq")
            pfile = next(x for x in subject_paths if x.name == chosen_file)

            notes_text = normalize_spaces(extract_text_smart(pfile))
            if not notes_text:
                st.warning("No text extracted (PDF may be scanned; OCR required).")
            else:
                if "quiz" not in st.session_state:
                    st.session_state.quiz = None
                if "quiz_subject" not in st.session_state:
                    st.session_state.quiz_subject = ""
                if "user_answers" not in st.session_state:
                    st.session_state.user_answers = {}

                st.caption("Step 1: Generate quiz. Step 2: Answer. Step 3: AI checks your answers.")
                n_items = st.slider("Number of questions", 5, 20, 10, key="nq")
                diff = st.select_slider("Difficulty", ["easy", "medium", "hard"], value="medium", key="qd")

                if st.button("🧩 Generate Quiz", key="gen_quiz"):
                    notes_limited = notes_text[:5500]
                    prompt = f"""
Return ONLY JSON. No markdown. No extra text.

Schema:
{{
  "questions":[
    {{
      "id":"Q1",
      "type":"mcq|identification|fill_blanks|short_answer",
      "question":"...",
      "choices":{{"A":"...","B":"...","C":"...","D":"..."}} (only for mcq),
      "answer_key":"A" (for mcq) OR "expected answer text" (for others)
    }}
  ]
}}

Rules:
- Use ONLY the notes.
- Mix types: include at least 3 different types.
- For fill_blanks: question must contain exactly ONE blank: "_____" and answer_key is the missing word/phrase.
- For short_answer/identification: answer_key should be a short correct answer (1–2 sentences max).
- Make {n_items} questions total.
- Grade level: Grade 12 STEM.
Subject: {assess_subject}
Difficulty: {diff}

NOTES:
{notes_limited}
"""
                    chat = client.chat.completions.create(
                        model="llama-3.3-70b-versatile",
                        messages=[
                            {"role": "system", "content": "You output valid JSON only."},
                            {"role": "user", "content": prompt},
                        ],
                        temperature=0.0,
                    )
                    raw = (chat.choices[0].message.content or "").strip()
                    quiz = parse_json_loose(raw)
                    st.session_state.quiz = quiz
                    st.session_state.quiz_subject = assess_subject
                    st.session_state.user_answers = {}
                    st.success("Quiz generated!")

                if st.session_state.quiz and st.session_state.quiz_subject == assess_subject:
                    quiz = st.session_state.quiz
                    questions = quiz.get("questions", [])

                    st.divider()
                    st.subheader("✍️ Answer the quiz")

                    for q in questions:
                        qid = q.get("id", "")
                        qtype = q.get("type", "")
                        qtext = q.get("question", "")
                        st.markdown(f"**{qid} ({qtype})**")
                        st.write(qtext)

                        if qtype == "mcq":
                            choices = q.get("choices", {})
                            options = [f"{k}) {choices[k]}" for k in ["A","B","C","D"] if k in choices]
                            pick = st.radio("Your answer", options, key=f"ans_{qid}", index=0 if options else 0)
                            st.session_state.user_answers[qid] = pick[:1] if pick else ""
                        else:
                            ans = st.text_area("Your answer", key=f"ans_{qid}", height=90)
                            st.session_state.user_answers[qid] = ans.strip()

                        st.write("---")

                    if st.button("✅ Check my answers (AI)", key="grade"):
                        payload = []
                        for q in questions:
                            qid = q.get("id", "")
                            payload.append(
                                {
                                    "id": qid,
                                    "type": q.get("type", ""),
                                    "question": q.get("question", ""),
                                    "choices": q.get("choices", {}),
                                    "expected": q.get("answer_key", ""),
                                    "user_answer": st.session_state.user_answers.get(qid, ""),
                                }
                            )

                        grade_prompt = f"""
Return ONLY JSON.

Schema:
{{
  "total_score": <number>,
  "max_score": <number>,
  "results": [
    {{
      "id":"Q1",
      "score": <0..1 for mcq, 0..2 for others>,
      "max_score": <number>,
      "correct": true/false,
      "feedback":"short helpful feedback",
      "suggestion":"what to review"
    }}
  ]
}}

Rules:
- For MCQ: max_score=1, correct only if user_answer matches expected letter.
- For non-MCQ: max_score=2, grade for correctness + completeness.
- Be strict but fair. Feedback must be short and useful.
- Use expected answers and notes context. Do not invent facts beyond notes.

Subject: {assess_subject}

DATA:
{json.dumps(payload, ensure_ascii=False)[:12000]}
"""
                        chat2 = client.chat.completions.create(
                            model="llama-3.3-70b-versatile",
                            messages=[
                                {"role": "system", "content": "You output valid JSON only."},
                                {"role": "user", "content": grade_prompt},
                            ],
                            temperature=0.0,
                        )
                        raw2 = (chat2.choices[0].message.content or "").strip()
                        graded = parse_json_loose(raw2)

                        st.subheader("📊 Results")
                        st.write(f"Score: **{graded.get('total_score')} / {graded.get('max_score')}**")

                        res = graded.get("results", [])
                        if isinstance(res, list) and res:
                            out_df = pd.DataFrame(res)
                            st.dataframe(out_df, use_container_width=True, hide_index=True)


# -----------------------------
# TAB: Dashboard
# -----------------------------
with tabs[4]:
    st.subheader("📊 Dashboard")

    if df.empty:
        st.info("Add subjects first.")
    else:
        most_urgent = df.sort_values("Days Left").iloc[0]["Subject"]
        next_exam = df["Exam Date"].min()

        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Minutes/day", f"{minutes_per_day}")
        c2.metric("Weekly minutes", f"{weekly_minutes}")
        c3.metric("Most urgent", str(most_urgent))
        c4.metric("Next exam", nice_date(next_exam))

        st.divider()
        st.subheader("Results")
        show_cols = ["Subject", "Difficulty (1-5)", "Exam Date", "Days Left", "Minutes/Week (Suggested)", "Minutes Done (this week)"]
        st.dataframe(
            df[show_cols].sort_values("Minutes/Week (Suggested)", ascending=False),
            use_container_width=True,
            hide_index=True,
            column_config={"Exam Date": st.column_config.DateColumn("Exam Date", format="MMM DD, YYYY")},
        )

        st.subheader("Progress")
        for _, row in df.sort_values("Minutes/Week (Suggested)", ascending=False).iterrows():
            done = int(row["Minutes Done (this week)"])
            target = int(row["Minutes/Week (Suggested)"])
            st.write(f"**{row['Subject']}** — {done}/{target} min")
            st.progress(float(row["Progress %"]))

        st.subheader("Study Plan")
        ranked = df.sort_values("Minutes/Week (Suggested)", ascending=False).reset_index(drop=True)
        top = ranked.head(min(3, len(ranked))).copy()

        if len(top) == 1:
            split = [1.0]
        elif len(top) == 2:
            split = [0.6, 0.4]
        else:
            split = [0.5, 0.3, 0.2]

        plan_rows = []
        for i in range(days_to_plan):
            day_label = (pd.Timestamp(date.today()) + pd.Timedelta(days=i)).strftime("%b %d, %Y")
            for idx in range(len(top)):
                plan_rows.append(
                    {"Day": day_label, "Subject": top.loc[idx, "Subject"], "Minutes": int(round(minutes_per_day * split[idx]))}
                )

        st.dataframe(pd.DataFrame(plan_rows), use_container_width=True, hide_index=True)

st.caption("✅ Keep-awake ping: /?ping=1")
